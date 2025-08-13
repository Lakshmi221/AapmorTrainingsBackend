from flask import Blueprint,request,jsonify
from werkzeug.utils import secure_filename
from datetime import datetime, timezone
import cloudinary.uploader
from pptx import Presentation
from io import BytesIO
from pymongo import ReturnDocument
from db import db
import requests
from bson import ObjectId
from utils.groq_slide_parser import classify_ppt
import fitz 
import os
from pathlib import Path
import traceback,inspect
from bson.errors import InvalidId
from utils.error_log import log_error_to_db
from utils.extract_images import extract_slide_images
from utils.compute_due_date import compute_due_date
import json
from flask import current_app
from utils.logger import setup_logger
logger = setup_logger("trainings_route")


trainings_bp=Blueprint('trainings',__name__)
file_path = os.path.relpath(inspect.getfile(inspect.currentframe()))

# upload
@trainings_bp.route('/trainings', methods=['POST'])
def create_training():
    logger.info("POST /trainings - create_training() called.")

    if 'file' not in request.files:
        logger.warning("File not provided in request.")
        error_msg = "File not provided"
        return jsonify({'error': error_msg}), 400
    
    file = request.files['file']

    if file.filename == '':
        logger.warning("Uploaded file has no filename.")
        return jsonify({'error':"file name is empty"}),404

    data = request.form
    # print(data)

    try:
        # upload_result = cloudinary.uploader.upload(
        #     file,
        #     resource_type='raw',
        #     folder='trainings_ppt',
        #     public_id=secure_filename(file.filename).split('.')[0]
        # )
        # ppt_url = upload_result.get('secure_url')
        file.stream.seek(0)
        ppt_stream=BytesIO(file.read())
        prs=Presentation(ppt_stream)
        

        if data.get('training_status') not in ['active', 'inactive']:
            return jsonify({"error": "trainings status must be 'active' or 'inactive'"}), 400
        if data.get("priority") not in ["high", "medium", "low"]:
            return jsonify({"error": "priority must be 'high', 'medium', or 'low'"}), 400

        enrolled_employees = data.get('enrolled_employees',0)

        # print(enrolled_employees)
        try:
            enrolled_employees = int(enrolled_employees)
            
        except (TypeError, ValueError):
            return jsonify({"error": "enrolled_employees must be a number"}), 400
        time_period = data.get('time_period')

        try:
            time_period = int(time_period)
        except (TypeError, ValueError):
            return jsonify({"error": "time_period must be a number"}), 400        

        raw_slides = []
        for idx, slide in enumerate(prs.slides):
            text_content = []
            images = []

            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.has_text_frame:
                    text_content.append(shape.text.strip())
                elif shape.shape_type == 13 and hasattr(shape, "image"):
                    image_blob = shape.image.blob
                    # print('true')
                    buffer = BytesIO(image_blob)
                    upload_result = cloudinary.uploader.upload(
                        buffer,
                        resource_type="image",
                        folder="training_slide_images",
                        public_id=f"slide_{idx+1}_{shape.shape_id}",
                        overwrite=True
                    )
                    images.append(upload_result["secure_url"])
                    

            raw_slides.append({
                'slide_number': idx + 1,
                'content': text_content,
                'images': images  
            })


        logger.info("Sending slide data to Groq for classification.")
        structured_blocks_raw = classify_ppt(raw_slides)
        
        if structured_blocks_raw:
            logger.debug(f"Groq response (first 100 chars): {structured_blocks_raw[:100]}")
        else:
            logger.error("No response received from Groq (structured_blocks_raw is None)")
            return jsonify({"error": "Failed to classify slides with Groq."}), 500


        
        try:
            structured_blocks = json.loads(structured_blocks_raw)
        except json.JSONDecodeError as err:
            logger.exception("Failed to parse Groq response as JSON.")
            return jsonify({
                "error": "Invalid JSON response from Groq",
                "detail": str(err),
                "raw_response": structured_blocks_raw
            }), 400

        
        training_title = structured_blocks[0]['title'] if structured_blocks[0]['type'] == 'title' else 'Untitled Training'
       
        training = {
            # 'ppt_url': ppt_url,
            
            'title':training_title,
            'ppt_data': structured_blocks,
            'uploaded_by_name': data.get('uploaded_by_name'),
            'uploaded_by_id': data.get('uploaded_by_id'),
            'uploaded_department': data.get('uploaded_department'),
            'timestamp': datetime.now(timezone.utc),
            'role': data.get('role'),
            'department': data.get('department'),
            'time_period': time_period,
            'enrolled_employees': enrolled_employees,
            'assigned_to':[],
            'trainings_status': data.get('training_status'),
            'priority': data.get('priority'),
            "completed":0,
            "due":0
            
        }
        


        result = db['trainings'].insert_one(training)
        logger.info(f"Training created successfully with ID: {str(result.inserted_id)}")


        return_training={
            'title':training_title,
            'time_period': time_period,
            'status':data.get('training_status'),
            'id': str(result.inserted_id),
            'due':0,
            'enrolled':0,
            'department':data.get('department'),
            'created_at':datetime.now(timezone.utc),
            "completed":0,     
        }
    
        return jsonify({"msg": "Training created", "id": str(result.inserted_id),"training":return_training}), 201
    except Exception as e:
        
        log_error_to_db(
            level="error",
            message=str(e),
            stack_trace=traceback.format_exc(),
            time=datetime.now(timezone.utc),
            path=file_path,
            method="create_training",
            ip=request.remote_addr
        )
        return jsonify({"error": "Something went wrong",'error':str(e)}), 500
    
#get all admin dashborad data
@trainings_bp.route('/admin-dashboard')
def get_admin_dashboard_data():
    try:
        response = requests.get('https://auth-x-api.azurewebsites.net/api/users')
        response.raise_for_status()
        employees = response.json()
    except requests.exceptions.RequestException as e:
        return jsonify({"error": "Failed to fetch employees", "details": str(e)}), 500

    trainings = list(db['trainings'].find())

    emp_map = {}
    for training in trainings:
        for assignment in training.get('assigned_to', []):
            emp_id = assignment.get('emp_id')
            if emp_id:
                if emp_id in emp_map:
                    emp_map[emp_id].append(assignment)
                else:
                    emp_map[emp_id] = [assignment] 

    employee_stats = []
    total_assigned = 0
    total_due = 0
    total_completed = 0
    
    # print(emp_map)

    for emp in employees:
        emp_id = emp.get('EmpId')
        emp_name = emp.get('FullName')
        emp_mail = emp.get('Email')
        emp_department = emp.get('SpaceName')

        assignments = emp_map.get(emp_id, [])
        # print(assignments)

        assigned = len(assignments)
        # due = sum(1 for a in assignments if str(a.get("status", "")).strip().lower() == "assigned")
        due = sum(1 for a in assignments if str(a.get("status", "")).strip().lower() in ["due", "assigned"])
        completed = sum(1 for a in assignments if str(a.get("status", "")).strip().lower() == "completed")
        
        # print("Assigned:", assigned, "Due:", due, "Completed:", completed)
        total_assigned += assigned
        total_due += due
        total_completed += completed

        employee_stats.append({
            'emp_id':emp_id,
            "name": emp_name,
            "email": emp_mail,
            "department": emp_department,
            "assignedTrainings": assigned,
            "dueTrainings": due,
            "completedTrainings": completed
        })

    trainings_list = []
    for training in trainings:
        trainings_list.append({
            "id": str(training.get("_id")),
            "title": training.get("title", "Untitled"),
            "department": training.get("department", "Unknown"),
            "status": "Active" if str(training.get("trainings_status", "")).lower() == "active" else "Inactive",
            'enrolled':training.get('enrolled_employees'),
            "time_period":training.get('time_period'),
            "created_at":training.get('timestamp'),
            "completed":training.get('completed'),
            'due':training.get("due")
        })

    dashboard_data = {
        "cardStats": {
            "assigned": total_assigned,
            "completed": total_completed,
            "due": total_due
        },
        "trainings": trainings_list,
        "employeeReport": employee_stats
    }

    return jsonify(dashboard_data)

@trainings_bp.route('/teamlead-dashboard', methods=['POST'])
def team_leadlead_dashboard():
    try:
        data = request.get_json()
        teamlead_email=data.get('email')
        team_dept=data.get('SpaceName')
        
        response = requests.get('https://auth-x-api.azurewebsites.net/api/users')
        response.raise_for_status()
        employees=response.json()
        
        team_members=[]
        
        
        
        for u in employees:
            if team_dept.lower() in u.get('SpaceName', '').lower() or u.get('SpaceName', '').lower() in team_dept.lower():
                team_members.append(u)
                # print(u)
        
        trainings=list(db['trainings'].find())
        
        team_stats = []
        total_assigned = 0
        total_completed = 0
        total_due = 0
        for member in team_members:
            emp_id=member.get('EmpId')
            # print(emp_id)
            name=member.get('FullName')
            role=member.get('SpaceName')
            email=member.get("Email")
            
            assigned_trainings = []
            assigned_training_ids = []
            for training in trainings:
                for a in training.get('assigned_to', []):
                    if str(a.get('emp_id')).strip() == str(emp_id).strip():
                        assigned_trainings.append(a)
                        assigned_training_ids.append(str(training.get('_id')))
            assigned = len(assigned_trainings)
            completed = sum(1 for a in assigned_trainings if a.get('status') == 'completed')
            due = sum(1 for a in assigned_trainings if str(a.get("status", "")).strip().lower() in ["due", "assigned"])
            progress = 0
            if assigned > 0:
                progress = int((completed / assigned) * 100)

            team_stats.append({
                "empId": emp_id,
                "name": name,
                "role": role,
                "email": email,
                "assigned": assigned,
                "completed": completed,
                "due": due,
                "assigned_trainings": assigned_training_ids
            })

            total_assigned += assigned
            total_completed += completed
            total_due += due   
            
        return jsonify({
            "cardStats": {
                "assigned": total_assigned,
                "completed": total_completed,
                "pending": total_due
            },
            "teamReport": team_stats
        })                         
        
    except Exception as e:
        return jsonify({"error": "Something went wrong", "details": str(e)}), 500
    
    
@trainings_bp.route('/trainings-for-teamlead', methods=['POST'])
def get_trainings_for_tl():
    try:
        data=request.get_json()
        department=data.get('department')
        all_trainings=list(db['trainings'].find())
        trainings=[]
        for training in all_trainings:
           if (training.get('department', '').lower() == department.lower()) or (training.get('department', '').lower() == 'all departments'):
                trainings.append({
                'title': training.get('title'),
                'id': str(training.get('_id')),
                'status':training.get('trainings_status')
            })
                
                
        return jsonify(trainings)
    except Exception as e:
        print(e)
        log_error_to_db(
            level="error",
            message=str(e),
            stack_trace=traceback.format_exc(),
            time=datetime.now(timezone.utc),
            path=file_path,
            method="get_training_by_id",
            ip=request.remote_addr
        )
        return jsonify({"error": "Failed to fetch"}), 500
         
    
#get single training
@trainings_bp.route('/trainings/<id>', methods=['GET'])
def get_training_by_id(id):
    try:
        emp_id = request.args.get("emp_id") 
        
        
        training = db['trainings'].find_one({"_id": ObjectId(id)})
        if not training:
            return jsonify({"error": "Training not found"}), 404
        
        assigned_list = training.get('assigned_to', [])
        user_entry = None
        for a in assigned_list:
            if a.get('emp_id')==emp_id:
                user_entry=a
                break
        training['assigned_to']=[user_entry]
        
        if not user_entry:
            return jsonify({"error": "Assignment not found for this employee"}), 404
        
        training['_id'] = str(training['_id'])
        return jsonify({'message': 'successfully Training fetched', 'trainings': training})
    except Exception as e:
        log_error_to_db(
            level="error",
            message=str(e),
            stack_trace=traceback.format_exc(),
            time=datetime.now(timezone.utc),
            path=file_path,
            method="get_training_by_id",
            ip=request.remote_addr
        )
        return jsonify({"error": "Failed to fetch"}), 500
    
# update
@trainings_bp.route('/update-trainings', methods=['PUT'])
def update_training_by_id():
    try:
        data = request.get_json()
        id=data.get('id')
        
        new_id=ObjectId(id)
        print(id) 
        change_obj=data.get('data')
        print(change_obj)

        training = db['trainings'].find_one({"_id": new_id})
        if not training:
            return jsonify({"error": "Training not found"}), 404

        
        update_fields = {}


        if 'training_status' in change_obj:
            status = change_obj['training_status']
            if status.lower() not in ['active', 'inactive']:
                return jsonify({'error': 'training status must be active/inactive'}), 400
            update_fields['trainings_status'] = status.lower()

    

        if change_obj.get('department') :
            update_fields['department'] = change_obj.get('department', "").strip()


        if change_obj.get('time_period'):
            time_period = change_obj.get('time_period')
            if time_period is not None:
                try:
                    update_fields['time_period'] = int(time_period)
                except (ValueError, TypeError):
                    return jsonify({"error": "time_period must be a number"}), 400


        if not update_fields:
            return jsonify({"error": "At least one field must be provided to update"}), 400

        update_fields['timestamp'] = datetime.now(timezone.utc)

        updated_training = db['trainings'].find_one_and_update(
            {"_id": ObjectId(id)},
            {"$set": update_fields},
            return_document=ReturnDocument.AFTER
        )

        return_doc={
            "id": str(updated_training.get('_id')),
            "title":updated_training.get("title"),
            "completed":updated_training.get('completed'),
            "due":updated_training.get('due'),
            "time_period":updated_training.get('time_period'),
            "status":updated_training.get('trainings_status'),
            "enrolled":updated_training.get("enrolled_employees"),
            "department":updated_training.get('department')
        }

        # print(updated_training)

        return jsonify({"msg": "Training updated", "trainings": return_doc})

    except Exception as e:
        print(e)
        log_error_to_db(
            level="error",
            message=str(e),
            stack_trace=traceback.format_exc(),
            time=datetime.now(timezone.utc),
            path=file_path,
            method="update_training_by_id",
            ip=request.remote_addr
        )
        return jsonify({"error": "Something went wrong", "details": str(e)}), 500



# delete
@trainings_bp.route('/trainings/<id>', methods=['DELETE'])
def delete_training_by_id(id):
    try:
        result = db['trainings'].delete_one({"_id": ObjectId(id)})

        if result.deleted_count == 0:
            msg = f"Delete failed: Training with ID '{id}' not found."
            current_app.logger.warning(msg)
            return jsonify({"error": "Training not found"}), 404

        current_app.logger.info(f"Training with ID '{id}' deleted successfully.")
        return jsonify({"msg": "Training deleted successfully","id":id}), 200

    except InvalidId:
        msg = f"Invalid ID format provided: {id}"
        current_app.logger.error(msg, exc_info=True)
        log_error_to_db(
            level="error",
            message=msg,
            stack_trace=traceback.format_exc(),
            time=datetime.now(timezone.utc),
            path=file_path,
            method="delete_training_by_id",
            ip=request.remote_addr
        )
        return jsonify({"error": "Invalid ID format"}), 400

    except Exception as e:
        msg = f"Unexpected error deleting training ID {id}"
        current_app.logger.error(msg, exc_info=True)
        log_error_to_db(
            level="error",
            message=msg,
            stack_trace=traceback.format_exc(),
            time=datetime.now(timezone.utc),
            path=file_path,
            method="delete_training_by_id",
            ip=request.remote_addr
        )
        return jsonify({"error": "Failed to delete training"}), 500

